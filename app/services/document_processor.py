"""
Document Processor Service
==========================
Comprehensive document processing supporting multiple formats:
- PDF (with text and image extraction + OCR)
- DOCX (Microsoft Word)
- PPTX (Microsoft PowerPoint)
- TXT, MD (Plain text and Markdown)
- ZIP (Archive containing multiple documents)

Optimized for:
- Memory efficiency (streaming where possible)
- Parallel processing for batch operations
- Selective OCR (only when text extraction fails)
"""

import os
import io
import re
import zipfile
import tempfile
import asyncio
import gc
import shutil
import logging
import importlib.util
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from datetime import datetime
import hashlib
from concurrent.futures import ThreadPoolExecutor

# Multimodal / Vision
try:
    import google.generativeai as genai

    VISION_AVAILABLE = True
except ImportError:
    genai = None
    VISION_AVAILABLE = False

# PDF Processing
from pypdf import PdfReader
import fitz  # PyMuPDF for image extraction

# Document Processing
from docx import Document as DocxDocument
from pptx import Presentation

# Image Processing
import numpy as np
from PIL import Image

# OCR (optional - graceful fallback if not available)
OCR_IMPORT_ERROR: Optional[str] = None
try:
    from paddleocr import PaddleOCR

    # Ensure paddle core dependency is present
    if importlib.util.find_spec("paddle") is None:
        raise ImportError("paddlepaddle is not installed")

    OCR_AVAILABLE = True
except ImportError as exc:
    PaddleOCR = None  # type: ignore
    OCR_AVAILABLE = False
    OCR_IMPORT_ERROR = str(exc)

from app.core.config import settings
from app.core.logging import get_logger
from app.services.vector_store import get_vector_store

logger = get_logger(__name__)

# Thread pool for CPU-bound tasks
_thread_pool = ThreadPoolExecutor(max_workers=2)


@dataclass
class ProcessedChunk:
    """A processed text chunk ready for embedding."""

    text: str
    metadata: Dict[str, Any]
    chunk_id: str


@dataclass
class ProcessedDocument:
    """Result of processing a single document."""

    filename: str
    file_type: str
    chunks: List[ProcessedChunk]
    page_count: int
    image_count: int
    total_characters: int
    processing_time_ms: float
    success: bool
    error: Optional[str] = None


@dataclass
class BatchProcessResult:
    """Result of batch processing multiple documents."""

    total_files: int
    successful_files: int
    failed_files: int
    total_chunks: int
    documents: List[ProcessedDocument]
    processing_time_ms: float


class DocumentProcessor:
    """
    Comprehensive document processor with support for multiple formats.

    Features:
    - Multi-format support (PDF, DOCX, PPTX, TXT, MD)
    - ZIP archive extraction and batch processing
    - PDF image extraction with OCR (optimized)
    - Smart text chunking with overlap
    - Metadata extraction
    - Memory-optimized processing
    - Content hash-based idempotency check

    Optimizations:
    - Uses temp files for large ZIPs instead of memory
    - Parallel processing for multiple files
    - Selective OCR only when text extraction yields little content
    - Batch vector store operations
    - Garbage collection for large files
    """

    # Supported file extensions
    SUPPORTED_EXTENSIONS = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "docx",  # Will attempt to process as docx
        ".pptx": "pptx",
        ".ppt": "pptx",  # Will attempt to process as pptx
        ".txt": "text",
        ".md": "markdown",
        ".markdown": "markdown",
        ".zip": "zip",
        ".jpg": "image",
        ".jpeg": "image",
        ".png": "image",
    }

    # OCR settings
    MIN_TEXT_LENGTH_FOR_OCR = 50  # Only OCR if page has less text than this
    MAX_IMAGES_PER_PAGE = 3  # Limit images to OCR per page
    MAX_IMAGE_SIZE = (1000, 1000)  # Resize large images before OCR

    # Parallel processing settings
    MAX_PARALLEL_FILES = 1  # Process sequentially to reduce peak memory

    def __init__(self):
        """Initialize document processor."""
        self.chunk_size = settings.CHUNK_SIZE
        self.chunk_overlap = settings.CHUNK_OVERLAP
        self.max_file_size = settings.MAX_FILE_SIZE_MB * 1024 * 1024
        self.max_zip_size = settings.MAX_ZIP_SIZE_MB * 1024 * 1024
        self.ocr_available = OCR_AVAILABLE and settings.ENABLE_OCR
        self.vision_available = (
            VISION_AVAILABLE and settings.ENABLE_MULTIMODAL_PROCESSING
        )
        self._ocr_engine = None
        self._vision_model = None

        # Idempotency: Track processed content hashes
        self._processed_hashes: Dict[str, str] = {}  # hash -> document_id

        if self.ocr_available:
            self._initialize_ocr_engine()

        if self.vision_available and settings.GEMINI_API_KEY:
            genai.configure(api_key=settings.GEMINI_API_KEY)
            self._vision_model = genai.GenerativeModel(settings.GEMINI_VISION_MODEL)
            logger.info("Gemini Vision initialized", model=settings.GEMINI_VISION_MODEL)
        elif settings.ENABLE_OCR:
            if OCR_IMPORT_ERROR:
                logger.warning(
                    "OCR disabled: PaddleOCR dependency unavailable",
                    reason=OCR_IMPORT_ERROR,
                )
            else:
                logger.warning("OCR enabled but PaddleOCR is not installed")
        else:
            logger.info("OCR disabled via configuration")

    def _compute_content_hash(self, content: bytes) -> str:
        """Compute SHA256 hash of file content for idempotency check."""
        return hashlib.sha256(content).hexdigest()

    def _check_duplicate(
        self, content_hash: str, collection_name: str
    ) -> Optional[str]:
        """
        Check if content has already been processed.

        Returns:
            document_id if duplicate, None otherwise
        """
        cache_key = f"{collection_name}:{content_hash}"
        return self._processed_hashes.get(cache_key)

    def _mark_processed(
        self, content_hash: str, collection_name: str, document_id: str
    ) -> None:
        """Mark content as processed for idempotency."""
        cache_key = f"{collection_name}:{content_hash}"
        self._processed_hashes[cache_key] = document_id

        # Prevent unbounded growth (keep last 10000 entries)
        if len(self._processed_hashes) > 10000:
            # Remove oldest entries (first 1000)
            keys_to_remove = list(self._processed_hashes.keys())[:1000]
            for key in keys_to_remove:
                del self._processed_hashes[key]

    def _compute_content_hash_from_path(self, file_path: str) -> str:
        """Compute SHA256 hash by streaming from disk (avoids loading entire file into RAM)."""
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    @staticmethod
    def _read_file_bytes(file_path: Optional[str]) -> bytes:
        """Read file from disk into bytes. Only used when streaming is not possible."""
        if file_path is None:
            raise ValueError("file_path is required when file_content is not provided")
        with open(file_path, "rb") as f:
            return f.read()

    async def process_file(
        self,
        file_content: Optional[bytes] = None,
        filename: str = "",
        document_id: str = "",
        collection_name: str = "",
        course_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
        file_path: Optional[str] = None,
    ) -> ProcessedDocument:
        """
        Process a single file of any supported type.

        Args:
            file_content: Raw file bytes (optional if file_path is provided)
            filename: Original filename
            document_id: Unique document identifier
            collection_name: Vector store collection name
            course_id: Optional course ID
            metadata: Optional additional metadata
            file_path: Optional path to file on disk (avoids loading into RAM)

        Returns:
            ProcessedDocument with chunks and stats
        """
        start_time = datetime.now()

        # If file_path is provided but no content, read only what's needed for hash
        if file_content is None and file_path is not None:
            # Compute hash by streaming the file instead of loading all into RAM
            content_hash = self._compute_content_hash_from_path(file_path)
            file_size = os.path.getsize(file_path)
        elif file_content is not None:
            content_hash = self._compute_content_hash(file_content)
            file_size = len(file_content)
        else:
            return ProcessedDocument(
                filename=filename,
                file_type="unknown",
                chunks=[],
                page_count=0,
                image_count=0,
                total_characters=0,
                processing_time_ms=0,
                success=False,
                error="Either file_content or file_path must be provided",
            )

        # Get file extension and type
        ext = Path(filename).suffix.lower()
        file_type = self.SUPPORTED_EXTENSIONS.get(ext)

        if not file_type:
            return ProcessedDocument(
                filename=filename,
                file_type="unknown",
                chunks=[],
                page_count=0,
                image_count=0,
                total_characters=0,
                processing_time_ms=0,
                success=False,
                error=f"Unsupported file type: {ext}",
            )

        # Idempotency check: Skip if identical content already processed
        existing_doc_id = self._check_duplicate(content_hash, collection_name)
        if existing_doc_id:
            logger.info(
                "document_already_processed",
                filename=filename,
                content_hash=content_hash[:16],
                existing_doc_id=existing_doc_id,
            )
            return ProcessedDocument(
                filename=filename,
                file_type=file_type,
                chunks=[],
                page_count=0,
                image_count=0,
                total_characters=0,
                processing_time_ms=0,
                success=True,
                error=f"Document already processed as {existing_doc_id}",
            )

        # Validate file size (use larger limit for ZIP files)
        max_size = self.max_zip_size if file_type == "zip" else self.max_file_size
        max_size_mb = (
            settings.MAX_ZIP_SIZE_MB
            if file_type == "zip"
            else settings.MAX_FILE_SIZE_MB
        )

        if file_size > max_size:
            return ProcessedDocument(
                filename=filename,
                file_type=file_type,
                chunks=[],
                page_count=0,
                image_count=0,
                total_characters=0,
                processing_time_ms=0,
                success=False,
                error=f"File size exceeds {max_size_mb}MB limit",
            )

        try:
            # Process based on file type
            if file_type == "pdf":
                result = await self._process_pdf(
                    file_content, filename, document_id, metadata, file_path=file_path
                )
            elif file_type == "docx":
                # DOCX/PPTX/text still need bytes; only load if not already provided
                content_for_processing = (
                    file_content
                    if file_content is not None
                    else self._read_file_bytes(file_path)
                )
                result = await self._process_docx(
                    content_for_processing, filename, document_id, metadata
                )
                if file_content is None:
                    del content_for_processing
                    gc.collect()
            elif file_type == "pptx":
                content_for_processing = (
                    file_content
                    if file_content is not None
                    else self._read_file_bytes(file_path)
                )
                result = await self._process_pptx(
                    content_for_processing, filename, document_id, metadata
                )
                if file_content is None:
                    del content_for_processing
                    gc.collect()
            elif file_type in ("text", "markdown"):
                content_for_processing = (
                    file_content
                    if file_content is not None
                    else self._read_file_bytes(file_path)
                )
                result = await self._process_text(
                    content_for_processing, filename, document_id, file_type, metadata
                )
                if file_content is None:
                    del content_for_processing
                    gc.collect()
            elif file_type == "zip":
                # ZIP processing returns multiple documents
                content_for_processing = (
                    file_content
                    if file_content is not None
                    else self._read_file_bytes(file_path)
                )
                batch_result = await self.process_zip(
                    content_for_processing,
                    document_id,
                    collection_name,
                    course_id,
                    metadata,
                )
                if file_content is None:
                    del content_for_processing
                    gc.collect()
                # Aggregate results
                total_chunks = sum(len(doc.chunks) for doc in batch_result.documents)
                return ProcessedDocument(
                    filename=filename,
                    file_type="zip",
                    chunks=[],  # Return empty list, they are already stored
                    page_count=batch_result.total_files,
                    image_count=0,
                    total_characters=sum(
                        doc.total_characters for doc in batch_result.documents
                    ),
                    processing_time_ms=batch_result.processing_time_ms,
                    success=batch_result.successful_files > 0,
                    error=None
                    if batch_result.successful_files > 0
                    else "No files processed successfully",
                )
            elif file_type == "image":
                content_for_processing = (
                    file_content
                    if file_content is not None
                    else self._read_file_bytes(file_path)
                )
                result = await self._process_image(
                    content_for_processing, filename, document_id, metadata
                )
                if file_content is None:
                    del content_for_processing
                    gc.collect()
            else:
                raise ValueError(f"Handler not implemented for type: {file_type}")

            # Store chunks in vector store
            if result.chunks:
                await self._store_chunks(result.chunks, collection_name)

            # Mark as processed for idempotency
            self._mark_processed(content_hash, collection_name, document_id)

            processing_time = (datetime.now() - start_time).total_seconds() * 1000
            result.processing_time_ms = processing_time

            logger.info(
                "document_processed",
                filename=filename,
                file_type=file_type,
                chunks=len(result.chunks),
                pages=result.page_count,
                images=result.image_count,
                processing_time_ms=processing_time,
                content_hash=content_hash[:16],
            )

            return result

        except Exception as e:
            processing_time = (datetime.now() - start_time).total_seconds() * 1000
            logger.error("document_processing_failed", filename=filename, error=str(e))

            return ProcessedDocument(
                filename=filename,
                file_type=file_type,
                chunks=[],
                page_count=0,
                image_count=0,
                total_characters=0,
                processing_time_ms=processing_time,
                success=False,
                error=str(e),
            )

    async def process_zip(
        self,
        zip_content: bytes,
        base_document_id: str,
        collection_name: str,
        course_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> BatchProcessResult:
        """
        Process a ZIP archive containing multiple documents.

        OPTIMIZED VERSION:
        - Extracts to temp directory (saves memory)
        - Processes files in parallel batches
        - Streams file content instead of loading all at once
        - Collects garbage after processing large files

        Args:
            zip_content: ZIP file bytes
            base_document_id: Base document ID (will be suffixed for each file)
            collection_name: Vector store collection name
            course_id: Optional course ID
            metadata: Optional additional metadata

        Returns:
            BatchProcessResult with all processed documents
        """
        start_time = datetime.now()
        documents: List[ProcessedDocument] = []
        temp_dir = None

        try:
            # Create temporary directory for extraction (saves memory vs in-memory processing)
            temp_dir = tempfile.mkdtemp(prefix="kolabri_zip_")
            zip_path = os.path.join(temp_dir, "archive.zip")

            # Write ZIP to temp file first
            with open(zip_path, "wb") as f:
                f.write(zip_content)

            # Free the zip_content from memory
            del zip_content
            gc.collect()

            # Extract files
            extract_dir = os.path.join(temp_dir, "extracted")
            os.makedirs(extract_dir, exist_ok=True)

            with zipfile.ZipFile(zip_path, "r") as zip_ref:
                # Security: Sanitize namelist to prevent Path Traversal (Zip Slip)
                file_list = []
                for f in zip_ref.namelist():
                    if (
                        f.endswith("/")
                        or Path(f).name.startswith(".")
                        or f.startswith("__MACOSX")
                    ):
                        continue

                    # Normalize path and check if it tries to go outside current dir
                    safe_path = os.path.normpath(f)
                    if safe_path.startswith("..") or os.path.isabs(safe_path):
                        logger.warning("skipping_unsafe_zip_path", path=f)
                        continue

                    if (
                        Path(f).suffix.lower() in self.SUPPORTED_EXTENSIONS
                        and Path(f).suffix.lower() != ".zip"
                    ):
                        file_list.append(f)

                logger.info("zip_extraction_started", total_files=len(file_list))
                zip_ref.extractall(extract_dir, members=file_list)

            # Remove the zip file to free space
            os.remove(zip_path)
            gc.collect()

            # Process files sequentially to keep peak memory low
            for global_index, file_path in enumerate(file_list):
                filename = Path(file_path).name
                full_path = os.path.join(extract_dir, file_path)

                if not os.path.exists(full_path):
                    continue

                logger.info(
                    "zip_processing_file",
                    filename=filename,
                    index=global_index + 1,
                    total=len(file_list),
                )

                result = await self._process_file_from_path(
                    file_path=full_path,
                    filename=filename,
                    document_id=f"{base_document_id}_{global_index}_{hashlib.md5(filename.encode()).hexdigest()[:8]}",
                    collection_name=collection_name,
                    course_id=course_id,
                    metadata={
                        **(metadata or {}),
                        "zip_source": True,
                        "original_path": file_path,
                    },
                )
                documents.append(result)

                # Garbage collect after each file to free decompressed data
                gc.collect()

            processing_time = (datetime.now() - start_time).total_seconds() * 1000

            successful = [d for d in documents if d.success]
            failed = [d for d in documents if not d.success]

            logger.info(
                "zip_processing_complete",
                total=len(documents),
                successful=len(successful),
                failed=len(failed),
                total_chunks=sum(len(d.chunks) for d in successful),
                processing_time_ms=processing_time,
            )

            return BatchProcessResult(
                total_files=len(documents),
                successful_files=len(successful),
                failed_files=len(failed),
                total_chunks=sum(len(d.chunks) for d in successful),
                documents=documents,
                processing_time_ms=processing_time,
            )

        except zipfile.BadZipFile:
            return BatchProcessResult(
                total_files=0,
                successful_files=0,
                failed_files=1,
                total_chunks=0,
                documents=[],
                processing_time_ms=(datetime.now() - start_time).total_seconds() * 1000,
            )
        finally:
            # Clean up temp directory
            if temp_dir and os.path.exists(temp_dir):
                try:
                    shutil.rmtree(temp_dir)
                except Exception as e:
                    logger.warning("temp_cleanup_failed", error=str(e))
            gc.collect()

    async def _process_file_from_path(
        self,
        file_path: str,
        filename: str,
        document_id: str,
        collection_name: str,
        course_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> ProcessedDocument:
        """
        Process a file from disk path (memory-efficient for large files).
        Passes file_path directly to process_file to avoid loading into RAM.
        """
        try:
            result = await self.process_file(
                file_content=None,
                filename=filename,
                document_id=document_id,
                collection_name=collection_name,
                course_id=course_id,
                metadata=metadata,
                file_path=file_path,
            )
            return result

        except Exception as e:
            logger.error("file_from_path_failed", path=file_path, error=str(e))
            return ProcessedDocument(
                filename=filename,
                file_type="unknown",
                chunks=[],
                page_count=0,
                image_count=0,
                total_characters=0,
                processing_time_ms=0,
                success=False,
                error=str(e),
            )
        finally:
            gc.collect()

    async def _process_pdf(
        self,
        content: Optional[bytes],
        filename: str,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None,
        file_path: Optional[str] = None,
    ) -> ProcessedDocument:
        """
        Process PDF with text and SELECTIVE image extraction.

        OPTIMIZED:
        - Opens PDF from disk path when available (avoids RAM duplication)
        - Only runs OCR if page text is below threshold
        - Limits number of images processed per page
        - Resizes large images before OCR
        - Uses thread pool for OCR operations
        """
        chunks: List[ProcessedChunk] = []
        all_text = []
        image_count = 0

        # Open PDF: prefer file_path (disk) over stream (RAM)
        if file_path is not None:
            pdf_doc = fitz.open(filename=file_path)
        elif content is not None:
            pdf_doc = fitz.open(stream=content, filetype="pdf")
        else:
            raise ValueError(
                "Either content or file_path must be provided for PDF processing"
            )
        page_count = len(pdf_doc)

        for page_num, page in enumerate(pdf_doc, start=1):
            page_text_parts = []

            # Extract text from page
            text = page.get_text("text")
            text_length = len(text.strip()) if text else 0

            if text.strip():
                page_text_parts.append(text.strip())

            # SELECTIVE OCR: Only if text extraction yielded little content
            should_ocr = (
                self.ocr_available and text_length < self.MIN_TEXT_LENGTH_FOR_OCR
            )

            if should_ocr:
                ocr_text = await self._run_page_ocr(page)
                if ocr_text:
                    page_text_parts.append(f"[OCR]: {ocr_text.strip()}")
                    image_count += 1

            # [NEW] MULTIMODAL: Extract images and generate captions
            if self.vision_available:
                image_list = page.get_images(full=True)
                for img_index, img in enumerate(image_list[: self.MAX_IMAGES_PER_PAGE]):
                    pil_img = None
                    try:
                        xref = img[0]
                        base_image = pdf_doc.extract_image(xref)
                        image_bytes = base_image["image"]

                        pil_img = Image.open(io.BytesIO(image_bytes))
                        if (
                            pil_img.width < settings.MIN_IMAGE_WIDTH
                            or pil_img.height < settings.MIN_IMAGE_HEIGHT
                        ):
                            pil_img.close()
                            del pil_img, image_bytes, base_image
                            pil_img = None
                            continue

                        logger.info(
                            "analyzing_pdf_image", page=page_num, img_index=img_index
                        )
                        caption = await self._generate_image_caption(pil_img)

                        if caption:
                            formatted_caption = (
                                f"\n\n=== [GAMBAR VISUAL DI HALAMAN {page_num}] ===\n"
                                f"Deskripsi: {caption}\n"
                                f"=============================================\\n\n"
                            )
                            page_text_parts.append(formatted_caption)
                            image_count += 1
                    except Exception as e:
                        logger.warning(
                            "pdf_image_extraction_failed", page=page_num, error=str(e)
                        )
                    finally:
                        # Aggressively free image memory per iteration
                        if pil_img is not None:
                            pil_img.close()
                            del pil_img
                        gc.collect()

            # Combine all text from page
            page_text = "\n\n".join(page_text_parts)
            if page_text:
                all_text.append(page_text)

                # Create chunks for this page
                page_chunks = self._create_chunks(
                    text=page_text,
                    document_id=document_id,
                    filename=filename,
                    page_number=page_num,
                    metadata=metadata,
                )
                chunks.extend(page_chunks)

        pdf_doc.close()
        gc.collect()

        return ProcessedDocument(
            filename=filename,
            file_type="pdf",
            chunks=chunks,
            page_count=page_count,
            image_count=image_count,
            total_characters=sum(len(c.text) for c in chunks),
            processing_time_ms=0,
            success=True,
        )

    async def _process_docx(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> ProcessedDocument:
        """Process DOCX (Microsoft Word) document."""
        doc = DocxDocument(io.BytesIO(content))

        # Extract all paragraphs
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)

        full_text = "\n\n".join(paragraphs)

        # [NEW] Extract images from DOCX — one at a time (generator)
        image_count = 0
        try:
            ocr_texts = []
            processed = 0
            for img in self._extract_images_from_docx(content):
                if processed >= 5:
                    img.close()
                    break
                try:
                    text = await self._run_ocr_optimized(img)
                    if text:
                        ocr_texts.append(f"[IMAGE_OCR]: {text}")
                finally:
                    img.close()
                    del img
                    gc.collect()
                processed += 1
            if ocr_texts:
                full_text += "\n\n" + "\n\n".join(ocr_texts)
                image_count = len(ocr_texts)
        except Exception as e:
            logger.warning("docx_image_extraction_failed", error=str(e))

        # Create chunks
        chunks = self._create_chunks(
            text=full_text,
            document_id=document_id,
            filename=filename,
            page_number=1,  # DOCX doesn't have pages in the same way
            metadata=metadata,
        )

        return ProcessedDocument(
            filename=filename,
            file_type="docx",
            chunks=chunks,
            page_count=1,
            image_count=image_count,
            total_characters=sum(len(c.text) for c in chunks),
            processing_time_ms=0,
            success=True,
        )

    async def _process_pptx(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> ProcessedDocument:
        """Process PPTX (Microsoft PowerPoint) presentation."""
        prs = Presentation(io.BytesIO(content))

        chunks: List[ProcessedChunk] = []
        total_chars = 0
        image_count = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())

                # Extract from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_text_parts.append(row_text)

                # [NEW] Extract and OCR images from PPTX shapes
                if shape.shape_type == 13:  # Picture
                    img = None
                    try:
                        image_data = shape.image.blob
                        img = Image.open(io.BytesIO(image_data))
                        del image_data
                        ocr_text = await self._run_ocr_optimized(img)
                        if ocr_text:
                            slide_text_parts.append(f"[IMAGE_OCR]: {ocr_text}")
                            image_count += 1
                    except Exception as e:
                        logger.warning("pptx_image_ocr_failed", error=str(e))
                    finally:
                        if img is not None:
                            img.close()
                            del img
                        gc.collect()

            slide_text = "\n".join(slide_text_parts)
            if slide_text:
                total_chars += len(slide_text)

                # Create chunks for this slide
                slide_chunks = self._create_chunks(
                    text=slide_text,
                    document_id=document_id,
                    filename=filename,
                    page_number=slide_num,
                    metadata={**(metadata or {}), "slide_number": slide_num},
                )
                chunks.extend(slide_chunks)

        return ProcessedDocument(
            filename=filename,
            file_type="pptx",
            chunks=chunks,
            page_count=len(prs.slides),
            image_count=image_count,
            total_characters=total_chars,
            processing_time_ms=0,
            success=True,
        )

    async def _process_text(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        file_type: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> ProcessedDocument:
        """Process plain text or markdown files."""
        # Try different encodings
        text = None
        for encoding in ["utf-8", "utf-16", "latin-1", "cp1252"]:
            try:
                text = content.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        if text is None:
            raise ValueError("Unable to decode text file with supported encodings")

        # Create chunks
        chunks = self._create_chunks(
            text=text,
            document_id=document_id,
            filename=filename,
            page_number=1,
            metadata=metadata,
        )

        return ProcessedDocument(
            filename=filename,
            file_type=file_type,
            chunks=chunks,
            page_count=1,
            image_count=0,
            total_characters=len(text),
            processing_time_ms=0,
            success=True,
        )

    async def _process_image(
        self,
        content: bytes,
        filename: str,
        document_id: str,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> ProcessedDocument:
        """Process a raw image file using Gemini Vision."""
        if not self.vision_available:
            return ProcessedDocument(
                filename=filename,
                file_type="image",
                chunks=[],
                page_count=0,
                image_count=0,
                total_characters=0,
                processing_time_ms=0,
                success=False,
                error="Multimodal/Vision processing is disabled or unavailable",
            )

        img = None
        try:
            img = Image.open(io.BytesIO(content))
            caption = await self._generate_image_caption(img)

            if not caption:
                raise ValueError("Vision AI failed to generate caption for image")

            full_text = f"=== [GAMBAR: {filename}] ===\nDeskripsi Visual: {caption}\n========================"

            chunks = self._create_chunks(
                text=full_text,
                document_id=document_id,
                filename=filename,
                page_number=1,
                metadata={**(metadata or {}), "is_multimodal": True},
            )

            return ProcessedDocument(
                filename=filename,
                file_type="image",
                chunks=chunks,
                page_count=1,
                image_count=1,
                total_characters=len(full_text),
                processing_time_ms=0,
                success=True,
            )
        except Exception as e:
            logger.error("image_processing_failed", filename=filename, error=str(e))
            raise
        finally:
            if img is not None:
                img.close()
                del img
            gc.collect()

    async def _generate_image_caption(self, image: Image.Image) -> str:
        """Generate description for an image using Gemini Vision."""
        if not self._vision_model:
            return ""

        try:
            prompt = """
            Analisis gambar ini secara detail untuk keperluan materi kuliah.
            1. Jika ini DIAGRAM/SKEMA: Jelaskan alur dan komponennya.
            2. Jika ini GRAFIK: Jelaskan sumbu X/Y, tren, dan titik penting.
            3. Jika ini RUMUS: Tuliskan dalam format LaTeX.
            4. Abaikan jika gambar buram atau tidak bermakna.
            
            Outputkan hanya deskripsinya saja dalam Bahasa Indonesia.
            """

            # Convert to JPEG bytes to release PIL image reference during API call
            img_buffer = io.BytesIO()
            rgb_image = image.convert("RGB") if image.mode != "RGB" else image
            rgb_image.save(img_buffer, format="JPEG", quality=85)
            img_bytes = img_buffer.getvalue()
            img_buffer.close()
            if rgb_image is not image:
                rgb_image.close()
                del rgb_image

            # Use bytes for API call so PIL image can be freed by caller
            import google.generativeai as genai_module

            img_part = {"mime_type": "image/jpeg", "data": img_bytes}

            loop = asyncio.get_running_loop()
            response = await loop.run_in_executor(
                _thread_pool,
                lambda: self._vision_model.generate_content([prompt, img_part]),
            )

            del img_bytes
            return response.text.strip()
        except Exception as e:
            logger.warning("vision_api_failed", error=str(e))
            return ""

    def _extract_images_from_docx(self, doc_content: bytes):
        """Extract images from DOCX bytes one at a time (generator to save RAM)."""
        try:
            with zipfile.ZipFile(io.BytesIO(doc_content)) as doc_zip:
                for name in doc_zip.namelist():
                    if name.startswith("word/media/"):
                        try:
                            image_data = doc_zip.read(name)
                            img = Image.open(io.BytesIO(image_data))
                            del image_data
                            yield img
                        except Exception:
                            continue
        except Exception as e:
            logger.warning("docx_extract_media_failed", error=str(e))

    async def _process_extracted_images(self, images: List[Image.Image]) -> str:
        """Process multiple images with OCR and return combined text."""
        combined_text = []
        for img in images:
            try:
                text = await self._run_ocr_optimized(img)
                if text:
                    combined_text.append(f"[IMAGE_OCR]: {text}")
            finally:
                img.close()
                del img
                gc.collect()
        return "\n\n".join(combined_text)

    def _initialize_ocr_engine(self) -> None:
        """Instantiate PaddleOCR once to avoid repeated heavyweight setup."""
        if self._ocr_engine is not None or not OCR_AVAILABLE:
            return

        try:
            lang = getattr(settings, "OCR_LANGUAGE", None) or "en"
            self._ocr_engine = PaddleOCR(
                use_angle_cls=True,
                lang=lang,
            )
            logging.getLogger("ppocr").setLevel(logging.ERROR)
            logger.info("PaddleOCR initialized", lang=lang)
        except Exception as exc:
            self._ocr_engine = None
            self.ocr_available = False
            logger.error("Failed to initialize PaddleOCR", error=str(exc))

    async def _run_ocr(self, image: Image.Image) -> str:
        """Run OCR on an image using PaddleOCR."""
        return await self._run_ocr_optimized(image)

    def _run_paddle_ocr(self, image: Image.Image) -> str:
        """Synchronous helper that executes PaddleOCR on a prepared PIL image."""
        if not self._ocr_engine:
            self._initialize_ocr_engine()

        engine = self._ocr_engine
        if engine is None:
            return ""

        try:
            np_image = np.array(image)
            result = engine.ocr(np_image, cls=True)
        except Exception as exc:
            logger.debug("ocr_failed", error=str(exc))
            return ""
        finally:
            # Free the numpy copy immediately
            try:
                del np_image
            except NameError:
                pass

        if not result:
            return ""

        lines: List[str] = []
        for line in result:
            for _, (text, confidence) in line:
                if not text:
                    continue
                if confidence is not None and confidence < 0.4:
                    continue
                lines.append(text.strip())

        return "\n".join(lines).strip()

    async def _run_ocr_optimized(self, image: Image.Image) -> str:
        """
        Run OCR on an image using PaddleOCR with memory and latency optimizations.

        Optimizations:
        - Resizes large images to reduce processing time
        - Uses thread pool to avoid blocking the event loop
        - Normalizes color space for consistent OCR results
        """
        if not self.ocr_available:
            return ""

        try:
            # Operate directly on the passed image to avoid RAM duplication
            if (
                image.width > self.MAX_IMAGE_SIZE[0]
                or image.height > self.MAX_IMAGE_SIZE[1]
            ):
                image.thumbnail(self.MAX_IMAGE_SIZE, Image.Resampling.LANCZOS)

            original_image = image
            if image.mode != "RGB":
                image = image.convert("RGB")
                # Close the original non-RGB image to free RAM
                original_image.close()
                del original_image

            loop = asyncio.get_running_loop()
            text = await loop.run_in_executor(
                _thread_pool, lambda: self._run_paddle_ocr(image)
            )

            return text.strip()
        except Exception as e:
            logger.debug("ocr_failed", error=str(e))
            return ""

    async def _run_page_ocr(self, page: fitz.Page) -> str:
        """Render a PDF page to image and run OCR (used when text extraction fails)."""
        if not self.ocr_available:
            return ""

        def render_page() -> Optional[Image.Image]:
            try:
                # Render page at ~112 DPI (reduced from 150 DPI to save RAM)
                matrix = fitz.Matrix(1.5, 1.5)
                pix = page.get_pixmap(matrix=matrix, alpha=False)
                mode = "RGB"
                image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
                pix = None
                return image
            except Exception as exc:
                logger.debug("page_render_failed", error=str(exc))
                return None

        loop = asyncio.get_running_loop()
        image = await loop.run_in_executor(_thread_pool, render_page)
        if image is None:
            return ""

        try:
            text = await self._run_ocr_optimized(image)
            return text
        finally:
            del image

    def _create_chunks(
        self,
        text: str,
        document_id: str,
        filename: str,
        page_number: int,
        metadata: Optional[Dict[str, Any]] = None,
    ) -> List[ProcessedChunk]:
        """Create overlapping chunks from text."""
        chunks = []

        # Clean text
        text = self._clean_text(text)

        if not text:
            return chunks

        # If text is smaller than chunk size, return as single chunk
        if len(text) <= self.chunk_size:
            chunk_id = f"{document_id}_p{page_number}_c0"
            chunks.append(
                ProcessedChunk(
                    text=text,
                    metadata={
                        **(metadata or {}),
                        "document_id": document_id,
                        "source": filename,
                        "page": page_number,
                        "chunk_index": 0,
                        "chunk_count": 1,
                    },
                    chunk_id=chunk_id,
                )
            )
            return chunks

        # Split into overlapping chunks
        start = 0
        chunk_index = 0

        while start < len(text):
            end = min(start + self.chunk_size, len(text))

            # Try to break at sentence or paragraph boundary
            if end < len(text):
                # Look for good break points
                for separator in ["\n\n", ". ", ".\n", "? ", "! ", "\n"]:
                    last_sep = text.rfind(separator, start + self.chunk_size // 2, end)
                    if last_sep > start:
                        end = last_sep + len(separator)
                        break

            chunk_text = text[start:end].strip()

            if chunk_text:
                chunk_id = f"{document_id}_p{page_number}_c{chunk_index}"
                chunks.append(
                    ProcessedChunk(
                        text=chunk_text,
                        metadata={
                            **(metadata or {}),
                            "document_id": document_id,
                            "source": filename,
                            "page": page_number,
                            "chunk_index": chunk_index,
                            "char_start": start,
                            "char_end": end,
                        },
                        chunk_id=chunk_id,
                    )
                )
                chunk_index += 1

            # Move to next position with overlap
            # Guard against infinite loop: start must always advance
            new_start = end - self.chunk_overlap
            if new_start <= start:
                # Overlap would go backwards or stay same — force advance past end
                start = end
            else:
                start = new_start
            if start >= len(text):
                break

        # Update chunk count in all metadata
        for chunk in chunks:
            chunk.metadata["chunk_count"] = len(chunks)

        return chunks

    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove excessive whitespace
        text = re.sub(r"\s+", " ", text)

        # Remove control characters except newlines
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)

        # Normalize newlines
        text = re.sub(r"\n{3,}", "\n\n", text)

        return text.strip()

    async def _store_chunks(
        self,
        chunks: List[ProcessedChunk],
        collection_name: str,
        batch_size: int = 100,
    ) -> None:
        """Store chunks in vector store with batching to limit memory usage."""
        if not chunks:
            return

        vector_store = get_vector_store()

        documents = [chunk.text for chunk in chunks]
        metadatas = [chunk.metadata for chunk in chunks]
        ids = [chunk.chunk_id for chunk in chunks]

        total = len(documents)
        for i in range(0, total, batch_size):
            batch_docs = documents[i : i + batch_size]
            batch_meta = metadatas[i : i + batch_size]
            batch_ids = ids[i : i + batch_size]

            await vector_store.add_documents(
                documents=batch_docs,
                metadatas=batch_meta,
                ids=batch_ids,
                collection_name=collection_name,
            )

            logger.info(
                "chunks_batch_stored",
                collection=collection_name,
                batch=f"{i // batch_size + 1}/{(total + batch_size - 1) // batch_size}",
                count=len(batch_docs),
            )

        logger.info("chunks_stored", collection=collection_name, count=total)


# Singleton instance
_document_processor: Optional[DocumentProcessor] = None


def get_document_processor() -> DocumentProcessor:
    """Get or create the document processor singleton."""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor
